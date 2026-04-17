from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    def add_slide(title_text, subtitle_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set dark background for forensic feel
        set_slide_background(slide, RGBColor(10, 10, 10))

        # Title styling
        title = slide.shapes.title
        title.text = title_text
        title_para = title.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(36)
        title_para.font.color.rgb = RGBColor(255, 77, 0) # VerifAI Orange
        title_para.alignment = PP_ALIGN.LEFT

        # Content styling
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = subtitle_text
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(220, 220, 220)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, RGBColor(10, 10, 10))
    title = slide.shapes.title
    title.text = "VerifAI – The Forensic Ledger"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 77, 0)
    subtitle = slide.placeholders[1]
    subtitle.text = "AI-Powered Multimodal Identity Verification & AML Scoring\nTechnical Overview"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    # Slide 2: Pipeline
    add_slide("The 6-Stage Forensic Pipeline", "A Modular Architecture for Absolute Integrity", [
        "Stage 1: ELA (Error Level Analysis) – Detecting digital tampering.",
        "Stage 2: AI Multimodal Extraction – Unified vision processing.",
        "Stage 3: Identity Mapping – Resolving name variations & typos.",
        "Stage 4: Compliance Engine – Regional rule validation (Regex/Freshness).",
        "Stage 5: [AI] Risk Analysis – Synthetic detection & PEP/AML scanning.",
        "Stage 6: [AI] Executive Summary – Automated forensic narration."
    ])

    # Slide 3: Vision Core
    add_slide("Multimodal Vision Reasoning (The AI Core)", "Leveraging GPT-4o for Cross-Document Intelligence", [
        "Unified Prompting: Aadhaar, PAN, and Utility Bills processed in a single vision context.",
        "Beyond OCR: The AI evaluates the relationship between documents (Cross-Document Mapping).",
        "Typo Resilience: Intelligent fuzzy matching allows for realistic spelling variations while maintaining security.",
        "Contextual Logic: Verifying addresses on utility bills against identity documents in real-time."
    ])

    # Slide 4: Synthetic Detection
    add_slide("AI-Powered Synthetic Detection", "Defending Against Deepfake & Generative Fraud", [
        "Signature Analysis: AI scans for 'too-perfect' text, unnatural lighting, and lack of physical grain.",
        "Synthetic Flags: Automatic detection of AI-generated documents triggers a high-risk veto.",
        "The Defense: Proactive identification of mathematically perfect fraudulent documents.",
        "Artifact Identification: Looking for photo-layer inconsistencies and background blur common in LLM generation."
    ])

    # Slide 5: Risk & AML
    add_slide("Predictive Risk & AML Scoring", "Intelligent Threat Assessment (0-100 Score)", [
        "AML Integration: Real-time fuzzy matching against the international PEP Ledger.",
        "Weighted Scoring: Risk score calculated from forensic signatures, AI certainty, and compliance pass-rates.",
        "Dynamic Output: Granular risk profiles instead of binary pass/fail results.",
        "Automated Veto: High-risk scores trigger manual escalation or immediate rejection."
    ])

    # Slide 6: Data Architecture
    add_slide("Zero-Knowledge Data Architecture", "Forensic Integrity & PII Privacy", [
        "Encryption-at-Rest: All sensitive PII (Names, IDs, AI Reasoning) is encrypted via Fernet.",
        "Immutable Ledger: Reports stored as encrypted documents in MongoDB Atlas for a tamper-proof audit trail.",
        "RBAC: Admin-only access to historical ledger data, secured by Google OAuth.",
        "Secure Storage: Local document persistence with UUID-based obfuscation."
    ])

    # Slide 7: Summary & Assistant
    add_slide("AI Executive Summary & Assistant", "Turning Raw Data into Forensic Intelligence", [
        "The Summary: AI synthesizes complex pipeline data into a concise narrative for the PDF report.",
        "VerifAI Assistant: A RAG-powered chatbot allowing natural language querying of the encrypted ledger.",
        "Natural Language Interface: Query example: 'Show me the forensic reasoning for the last rejected Aadhaar.'",
        "Secure Context: Assistant decrypts data in-memory only for authorized sessions."
    ])

    # Slide 8: Technical Impact
    add_slide("Technical Impact & Future-Proofing", "The New Standard in Forensic Verification", [
        "Scalability: Asynchronous Python architecture handles high-concurrency loads.",
        "Compliance: Built for GDPR and DPDP standards from the ground up.",
        "The Result: A high-fidelity, AI-driven forensic tool that identifies fraud where traditional systems fail.",
        "Extensibility: Strategy pattern allows for easy addition of new document types and forensic stages."
    ])

    prs.save('VerifAI_Forensic_Overview.pptx')
    print("✅ PowerPoint created successfully: VerifAI_Forensic_Overview.pptx")

if __name__ == "__main__":
    create_presentation()
